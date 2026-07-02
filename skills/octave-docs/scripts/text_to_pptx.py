#!/usr/bin/env python3
"""Convert a plain-text outline into a first-pass PPTX using the Octave POTX template.

Outline format (markdown-like):
  # Deck Title
  ## Section Name
  A section header slide (no bullets or table rows under it).
  ## Slide Title
  - bullet one
  - bullet two
  ## Table Slide
  | Header One | Header Two |
  | row one col a | row one col b |

Usage:
  python3 text_to_pptx.py <outline.txt> <output.pptx>
"""
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

ASSET = Path(__file__).parent.parent / "assets" / "PowerPoint" / "Octave_PPT_Template_20260401.potx"

TEMPLATE_CT = b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
PRESENTATION_CT = b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"

# Position/size of the table shape on the POTX's own "Title and Table" example slide
# (the layout has no table placeholder -- the template's own example adds it as a
# free-floating graphicFrame at this exact position).
TABLE_POSITION = (Emu(1323975), Emu(1380618))
TABLE_SIZE = (Emu(10487025), Emu(2287461))
OCTAVE_TABLE_STYLE_ID = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def load_template(potx_path: Path, tmp_pptx: Path) -> Presentation:
    """python-pptx refuses .potx directly (content type is 'template', not
    'presentation'). Patch [Content_Types].xml so it opens as a normal deck."""
    with zipfile.ZipFile(potx_path) as zin:
        with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                if name == "[Content_Types].xml":
                    data = data.replace(TEMPLATE_CT, PRESENTATION_CT)
                zout.writestr(name, data)
    return Presentation(tmp_pptx)


def strip_example_slides(prs: Presentation) -> None:
    """The POTX ships 75 pre-built example slides (one per layout) as a
    gallery. Every generated deck must start empty, or the user's content
    gets appended after the entire example gallery."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def _parse_table_row(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def parse_outline(text: str):
    """Returns (title, sections). Each section is (heading, kind, payload):
    kind "section_header" (payload None), "content" (payload list[str] bullets),
    or "table" (payload list[list[str]] rows, first row is the header row)."""
    lines = [ln.rstrip() for ln in text.splitlines() if ln.strip()]
    if not lines or not lines[0].startswith("# "):
        raise ValueError("Outline must start with '# Deck Title'")
    title = lines[0][2:].strip()

    def close_section(heading, bullets, rows):
        if rows:
            return (heading, "table", rows)
        if bullets:
            return (heading, "content", bullets)
        return (heading, "section_header", None)

    sections = []
    heading = None
    bullets: list[str] = []
    rows: list[list[str]] = []
    for line in lines[1:]:
        if line.startswith("## "):
            if heading is not None:
                sections.append(close_section(heading, bullets, rows))
            heading = line[3:].strip()
            bullets, rows = [], []
        elif line.startswith("- "):
            bullets.append(line[2:].strip())
        elif line.startswith("|"):
            rows.append(_parse_table_row(line))
    if heading is not None:
        sections.append(close_section(heading, bullets, rows))

    return title, sections


def _add_table_slide(prs: Presentation, layout, heading: str, rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(layout)
    slide.placeholders[0].text = heading

    n_rows, n_cols = len(rows), len(rows[0])
    left, top = TABLE_POSITION
    width, height = TABLE_SIZE
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            table.cell(r, c).text = cell_text

    # python-pptx has no public API for setting a specific registered table style,
    # so apply Octave's brand style directly via the underlying XML.
    tbl_pr = table._tbl.find(f"{A_NS}tblPr")
    style_id_el = tbl_pr.find(f"{A_NS}tableStyleId")
    if style_id_el is None:
        style_id_el = etree.SubElement(tbl_pr, f"{A_NS}tableStyleId")
    style_id_el.text = OCTAVE_TABLE_STYLE_ID


def build_pptx(outline_text: str, output_path: Path) -> None:
    tmp_pptx = output_path.with_suffix(".base.pptx")
    prs = load_template(ASSET, tmp_pptx)
    strip_example_slides(prs)

    layouts = {layout.name: layout for layout in prs.slide_masters[0].slide_layouts}
    title, sections = parse_outline(outline_text)

    title_slide = prs.slides.add_slide(layouts["Title Slide"])
    title_slide.placeholders[0].text = title

    for heading, kind, payload in sections:
        if kind == "section_header":
            slide = prs.slides.add_slide(layouts["Section Header"])
            slide.placeholders[0].text = heading
            continue

        if kind == "table":
            _add_table_slide(prs, layouts["Title and Table"], heading, payload)
            continue

        slide = prs.slides.add_slide(layouts["Title and Content"])
        slide.placeholders[0].text = heading
        body = slide.placeholders[1].text_frame
        body.text = payload[0]
        for bullet in payload[1:]:
            p = body.add_paragraph()
            p.text = bullet

    prs.save(output_path)
    tmp_pptx.unlink(missing_ok=True)


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: text_to_pptx.py <outline.txt> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    outline_path, out_path = Path(sys.argv[1]), Path(sys.argv[2])
    build_pptx(outline_path.read_text(), out_path)
    print(f"Wrote {out_path}")
